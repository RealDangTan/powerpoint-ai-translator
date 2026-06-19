import time
import json
import hashlib
import re
import os
import httpx
import random
from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ============================================================
# CẤU HÌNH
# ============================================================
ORIMISE_BASE_URL = "ENDPOINT_URL_HERE"  # Replace with your actual endpoint URL
ORIMISE_API_KEY = "YOUR_API_KEY_HERE"  # Replace with your actual API key

# Flags điều khiển demo
SHOW_USAGE = False
SHOW_GEMINI = False

# Sử dụng trực tiếp model xịn nhất
MODEL = "claude-sonnet-4-6"

CACHE_FILE = "translation_cache_v2.json"

INPUT_FILE = "input.pptx"
OUTPUT_FILE = "output.pptx"

TRANSLATION_DIRECTION = "VI_TO_EN"  # Options: "EN_TO_VI" or "VI_TO_EN"

# ============================================================
# SYSTEM PROMPTS
# ============================================================
CONTEXT_EXTRACTOR_PROMPT = """Role: Context Extractor.
Input: List of slide titles/headers.
Task: Summarize the main topic, industry, and tone (formal/creative/technical) in 20 words.
Output: Topic: [Target], Industry: [Target], Tone: [Target]."""

BATCH_TRANSLATOR_PROMPT_EN_TO_VI = """Role: Senior Translator.
Context: {global_context}
MANDATORY GLOSSARY REPLACEMENT (APPLY FIRST):
- Replace ALL occurrences of these English terms with their exact Vietnamese equivalents BEFORE translating:
{glossary_lines}
Task: Translate English → Professional Vietnamese.
Rules:
- Keep ID prefixes (T0:, T1:, etc.) intact.
- One translation per line, format: ID: translated_text
- No explanations. Direct translation only.
- Use industry-specific terminology from Context.
- If a term is a number or untranslatable, return as-is.
- If ambiguous, refer to Context.
- Maintain the same meaning and tone as the original English text, but adapt it to be natural and fluent in Vietnamese. Consider cultural nuances and local expressions.
- Don't translate word-by-word. Use context for fluent, natural Vietnamese — not literal translation."""

BATCH_TRANSLATOR_PROMPT_VI_TO_EN = """Role: Senior Translator.
Context: {global_context}
MANDATORY GLOSSARY REPLACEMENT (APPLY FIRST):
- Replace ALL occurrences of these Vietnamese terms with their exact English equivalents BEFORE translating:
{glossary_lines}
Task: Translate Vietnamese → Professional English.
Rules:
- Keep ID prefixes (T0:, T1:, etc.) intact.
- One translation per line, format: ID: translated_text
- No explanations. Direct translation only.
- Use industry-specific terminology from Context.
- If a term is a number or untranslatable, return as-is.
- If ambiguous, refer to Context.
- Maintain the same meaning and tone as the original Vietnamese text, but adapt it to be natural and fluent in English. Consider cultural nuances and local expressions.
- Don't translate word-by-word. Use context for fluent, natural English — not literal translation."""

# ============================================================
# GLOSSARY CHUYÊN NGÀNH
# ============================================================
MASTER_GLOSSARY = {
    "customer segment/base": "tệp khách hàng",
    "driving sales": "chạy số",
    "customer journey": "hành trình khách hàng",
    "customer experience": "trải nghiệm khách hàng",
    "touchpoint": "điểm chạm",
    "customer support": "CSKH",
    "persona": "chân dung",
    "QS": "CLDV",
    "CX": "TNKH",
    "employee": "CBNV",
    "relationship manager": "RM",
    "universal banker": "UB",
    "excellent CX": "trải nghiệm khách hàng xuất xắc",
    "listening post": "điểm lắng nghe",
    "P&S": "SPDV",
    "omnichannel": "đa kênh",
    "understand": "thấu hiểu",
    "listen": "lắng nghe",
    "act": "hành động",
    "best practice": "thông lệ",
    "closed-loop feedback": "phản hồi khách hàng",
    "engine": "nền tảng",
    "chương trình lắng nghe khách hàng": "Voice of Customer (VoC)",
}


# ============================================================
# HELPERS
# ============================================================
def _is_placeholder(val: str) -> bool:
    """Return True if the value looks like a placeholder."""
    if not val:
        return True
    return ("<" in val and ">" in val) or val.strip().startswith("<")


def _retry_call(fn, *args, retries=3, backoff=3.0, exceptions=(Exception,), **kwargs):
    """Retry wrapper with exponential backoff + jitter."""
    last_exc = None
    for attempt in range(1, retries + 1):
        try:
            return fn(*args, **kwargs)
        except exceptions as e:
            last_exc = e
            if attempt == retries:
                raise
            delay = backoff * (2 ** (attempt - 1)) + random.uniform(0, 1)
            print(f"      ⏳ Retry {attempt}/{retries} in {delay:.1f}s...")
            time.sleep(delay)
    raise last_exc


# ============================================================
# API CALLS (httpx only)
# ============================================================
def get_usage(base_url, api_key, start_date=None, end_date=None, params_extra=None):
    """Call the provider Usage endpoint via httpx."""
    url = f"{base_url.rstrip('/')}/v1/usage"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Accept": "application/json",
    }
    params = {}
    if start_date:
        params["start_date"] = start_date
    if end_date:
        params["end_date"] = end_date
    if params_extra:
        params.update(params_extra)

    resp = httpx.get(url, headers=headers, params=params, timeout=15)
    resp.raise_for_status()
    return resp.json()


def call_llm_generate(prompt_text, base_url=ORIMISE_BASE_URL,
                      api_key=ORIMISE_API_KEY, model="gpt-5.4"):
    """Call the API via OpenAI SDK."""
    if _is_placeholder(base_url) or _is_placeholder(api_key):
        raise RuntimeError("API base URL or API key is a placeholder.")

    # Append /v1 to base_url for OpenAI SDK compatibility
    client = OpenAI(base_url=f"{base_url.rstrip('/')}/v1", api_key=api_key)

    def _do_request():
        response = client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt_text}],
            timeout=30.0
        )
        return response.choices[0].message.content

    try:
        return _retry_call(_do_request, retries=3, backoff=50.0, exceptions=(Exception,))
    except Exception as e:
        print(f"   [DEBUG] Error calling {model}: {e}")
        raise


# ============================================================
# DEMOS
# ============================================================
def _maybe_print_usage_demo():
    """Opt-in usage demo."""
    if not SHOW_USAGE:
        return
    if _is_placeholder(ORIMISE_API_KEY):
        print("⚠️ SHOW_USAGE=True but API key is placeholder; skipping.")
        return
    print("\n🔎 Fetching usage...")
    try:
        usage = get_usage(ORIMISE_BASE_URL, ORIMISE_API_KEY)
        print(json.dumps(usage, indent=2, ensure_ascii=False))
    except Exception as e:
        print(f"❌ Usage fetch failed: {e}")


def _maybe_demo_gemini():
    """Opt-in Gemini demo."""
    if not SHOW_GEMINI:
        return
    if _is_placeholder(ORIMISE_API_KEY):
        print("⚠️ SHOW_GEMINI=True but API key is placeholder; skipping.")
        return
    try:
        print(f"\n🤖 LLM demo ({MODEL})...")
        out = call_llm_generate("Hello!")
        print(f"LLM says: {out}")
    except Exception as e:
        print(f"❌ Gemini demo failed: {e}")


# ============================================================
# CACHE
# ============================================================
def _make_cache_key(text, extra_context=""):
    raw = f"{text.strip()}||{extra_context}"
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()


def load_cache():
    if os.path.exists(CACHE_FILE):
        try:
            with open(CACHE_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                print(f"📦 Cache loaded: {len(data)} entries.")
                return data
        except (json.JSONDecodeError, IOError) as e:
            print(f"⚠️ Cache read error, creating new: {e}")
            return {}
    print("📦 No cache found, will create.")
    return {}


def save_cache(cache):
    try:
        with open(CACHE_FILE, "w", encoding="utf-8") as f:
            json.dump(cache, f, ensure_ascii=False, indent=2)
    except IOError as e:
        print(f"⚠️ Cache save error: {e}")


_stats = {"cache_hits": 0, "api_calls": 0, "tokens_saved_est": 0}


# ============================================================
# PPTX TEXT EXTRACTION
# ============================================================
def extract_all_text_from_shape(shape):
    """Recursively extract all text from a shape (including groups & tables)."""
    texts = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            texts.extend(extract_all_text_from_shape(sub))
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip()
            if t:
                texts.append(t)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        texts.append(t)
    return texts


def extract_slide_titles(prs):
    """Extract title from each slide."""
    titles = []
    for i, slide in enumerate(prs.slides):
        slide_title = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 0:
                t = shape.text.strip()
                if t:
                    slide_title = t
                    break
        if not slide_title:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        slide_title = t[:100]
                        break
        if slide_title:
            titles.append(f"Slide {i + 1}: {slide_title}")
    return titles


# ============================================================
# BƯỚC 1: GLOBAL CONTEXT EXTRACTION (single API call)
# ============================================================
def extract_global_context(prs, cache):
    """Extract overall context from presentation using a cheap model call."""
    titles = extract_slide_titles(prs)
    titles_text = "\n".join(titles)

    first_slide_content = ""
    if prs.slides:
        first_texts = []
        for shape in prs.slides[0].shapes:
            first_texts.extend(extract_all_text_from_shape(shape))
        first_slide_content = "\n".join(first_texts[:10])

    input_for_context = f"TITLES:\n{titles_text}\n\nFIRST SLIDE CONTENT:\n{first_slide_content}"

    cache_key = _make_cache_key(input_for_context, "global_context")
    if cache_key in cache:
        _stats["cache_hits"] += 1
        ctx = cache[cache_key]["result"]
        print(f"   ✅ Global Context (cached): {ctx}")
        return ctx

    time.sleep(2.0)  # Avoid duplicate detection

    try:
        prompt_text = f"SYSTEM:\n{CONTEXT_EXTRACTOR_PROMPT}\n\nUSER:\n{input_for_context}"

        ctx = call_llm_generate(
            prompt_text, base_url=ORIMISE_BASE_URL,
            api_key=ORIMISE_API_KEY, model=MODEL
        ).strip()

        cache[cache_key] = {
            "type": "global_context",
            "input_preview": input_for_context[:200],
            "result": ctx,
            "model": MODEL,
            "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
        }
        save_cache(cache)
        _stats["api_calls"] += 1
        print(f"   🌐 Global Context (API): {ctx}")
        return ctx

    except Exception as e:
        print(f"❌ Context extraction error: {e}")
        return "Topic: Business, Industry: General, Tone: Formal"


# ============================================================
# BƯỚC 2: PARTIAL GLOSSARY — only send relevant terms
# ============================================================
def build_partial_glossary(slide_text):
    slide_lower = slide_text.lower()
    relevant = [
        f"{en}:{vi}"
        for en, vi in MASTER_GLOSSARY.items()
        if en.lower() in slide_lower
    ]
    return "|".join(relevant) if relevant else ""


# ============================================================
# BƯỚC 3: BATCH TRANSLATION — one API call per slide
# ============================================================
def collect_slide_text_items(slide):
    """Collect all text paragraphs from a slide with IDs."""
    items = []
    counter = [0]

    def _collect_from_text_frame(text_frame):
        for paragraph in text_frame.paragraphs:
            full_text = "".join(run.text for run in paragraph.runs).strip()
            if full_text:
                items.append({
                    "id": f"T{counter[0]}",
                    "text": full_text,
                    "paragraph": paragraph,
                })
                counter[0] += 1

    def _collect_from_shape(shape):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                _collect_from_shape(sub)
        if shape.has_text_frame:
            _collect_from_text_frame(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    _collect_from_text_frame(cell.text_frame)

    for shape in slide.shapes:
        _collect_from_shape(shape)

    return items


def batch_translate_slide(items, global_context, cache):
    """Translate all text items in a slide via a single batched API call."""
    if not items:
        return

    time.sleep(1.5)  # Prevent duplicate request detection

    # Separate cached vs uncached
    cached_results = {}
    uncached_items = []
    all_text_combined = " ".join(item["text"] for item in items)

    for item in items:
        cache_key = _make_cache_key(
            item["text"], f"ctx:{global_context}|model:{MODEL}|dir:{TRANSLATION_DIRECTION}"
        )
        if cache_key in cache:
            cached_results[item["id"]] = cache[cache_key]["translated"]
            _stats["cache_hits"] += 1
            _stats["tokens_saved_est"] += len(item["text"].split()) * 3
        else:
            uncached_items.append(item)

    if cached_results:
        print(f"      ✅ {len(cached_results)} paragraphs from cache")

    # Translate uncached items
    if uncached_items:
        batch_lines = [f"{it['id']}: {it['text']}" for it in uncached_items]
        batch_input = "\n".join(batch_lines)

        if TRANSLATION_DIRECTION == "EN_TO_VI":
            glossary_lines = "\n".join(
                f"  - {en} → {vi}"
                for en, vi in MASTER_GLOSSARY.items()
                if en.lower() in all_text_combined.lower()
            )
            prompt_template = BATCH_TRANSLATOR_PROMPT_EN_TO_VI
        else:
            glossary_lines = "\n".join(
                f"  - {vi} → {en}"
                for en, vi in MASTER_GLOSSARY.items()
                if vi.lower() in all_text_combined.lower()
            )
            prompt_template = BATCH_TRANSLATOR_PROMPT_VI_TO_EN

        system_prompt = prompt_template.format(
            global_context=global_context,
            glossary_lines=glossary_lines
        )

        try:
            prompt_text = system_prompt + "\n\n" + batch_input
            raw_output = call_llm_generate(
                prompt_text, base_url=ORIMISE_BASE_URL,
                api_key=ORIMISE_API_KEY, model=MODEL
            ).strip()
            _stats["api_calls"] += 1

            api_results = _parse_batch_output(raw_output, uncached_items)

            for item in uncached_items:
                translated = api_results.get(item["id"], item["text"])
                cache_key = _make_cache_key(
                    item["text"], f"ctx:{global_context}|model:{MODEL}|dir:{TRANSLATION_DIRECTION}"
                )
                cache[cache_key] = {
                    "original": item["text"],
                    "translated": translated,
                    "context": global_context[:100],
                    "model": MODEL,
                    "timestamp": time.strftime("%Y-%m-%d %H:%M:%S"),
                }
                cached_results[item["id"]] = translated

            save_cache(cache)
            print(f"      🌐 {len(uncached_items)} paragraphs translated (1 API call)")

        except Exception as e:
            print(f"      ❌ Batch API error: {e}")
            for item in uncached_items:
                cached_results[item["id"]] = item["text"]

    # Write translations back to PPTX
    for item in items:
        translated = cached_results.get(item["id"], item["text"])
        paragraph = item["paragraph"]
        if paragraph.runs:
            paragraph.runs[0].text = translated
            for i in range(1, len(paragraph.runs)):
                paragraph.runs[i].text = ""


def _parse_batch_output(raw_output, uncached_items):
    """Parse batch output like 'T0: translated text'."""
    results = {}
    item_ids = {it["id"] for it in uncached_items}

    for line in raw_output.strip().split("\n"):
        line = line.strip()
        if not line:
            continue
        match = re.match(r"^(T\d+)\s*:\s*(.+)$", line)
        if match:
            item_id = match.group(1)
            translated = match.group(2).strip()
            if item_id in item_ids:
                results[item_id] = translated

    # Fill missing with originals
    if len(results) < len(uncached_items):
        missing = [it["id"] for it in uncached_items if it["id"] not in results]
        if missing:
            print(f"      ⚠️ Could not parse {len(missing)} IDs: {missing}")
            for it in uncached_items:
                if it["id"] not in results:
                    results[it["id"]] = it["text"]

    return results


# ============================================================
# MAIN
# ============================================================
def main(input_pptx, output_pptx):
    print("=" * 60)
    print(f"🚀 ADVANCED TRANSLATOR — {input_pptx}")
    print(f"   Model : {MODEL}")
    print("=" * 60)

    # Run demos if enabled
    _maybe_print_usage_demo()
    _maybe_demo_gemini()

    cache = load_cache()
    _stats["cache_hits"] = 0
    _stats["api_calls"] = 0
    _stats["tokens_saved_est"] = 0

    prs = Presentation(input_pptx)
    total_slides = len(prs.slides)

    # Step 1: Global context
    print(f"\n🔍 Step 1: Extracting global context...")
    global_context = extract_global_context(prs, cache)

    # Step 2: Translate slides
    print(f"\n📝 Step 2: Translating {total_slides} slides (batch mode)...")
    for i, slide in enumerate(prs.slides):
        print(f"\n   📄 Slide {i + 1}/{total_slides}")
        items = collect_slide_text_items(slide)
        if not items:
            print("      (empty, skipping)")
            continue
        print(f"      📋 {len(items)} text items")
        batch_translate_slide(items, global_context, cache)
        time.sleep(0.2)

    prs.save(output_pptx)

    # Stats
    total_ops = _stats["cache_hits"] + _stats["api_calls"]
    print(f"\n{'=' * 60}")
    print(f"✨ DONE! → {output_pptx}")
    print(f"{'=' * 60}")
    print("📊 PERFORMANCE:")
    print(f"   Total operations : {total_ops}")
    if total_ops > 0:
        pct = _stats["cache_hits"] / total_ops * 100
        print(f"   ✅ Cache hits    : {_stats['cache_hits']} ({pct:.1f}%)")
    print(f"   🌐 API calls     : {_stats['api_calls']}")
    print(f"   💾 Cache entries  : {len(cache)}")
    if _stats["tokens_saved_est"] > 0:
        print(f"   🪙 Tokens saved  : ~{_stats['tokens_saved_est']:,} (est.)")
    print(f"{'=' * 60}")


if __name__ == "__main__":
    main(INPUT_FILE, OUTPUT_FILE)